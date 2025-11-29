from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from io import BytesIO
import os

def create_presentation(slides_data, filename="generated_presentation.pptx"):
    """
    Creates a PowerPoint presentation from a list of slide data.
    Each slide data should have 'title', 'content' (list), and optionally 'image' (PIL Image).
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for slide_data in slides_data:
        # Use blank layout for complete control
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Add Title at the top
        title_box = slide.shapes.add_textbox(
            left=Inches(0.5),
            top=Inches(0.3),
            width=Inches(9),
            height=Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.text = slide_data.get("title", "Untitled")
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.LEFT

        # Determine if we have an image
        has_image = "image_bytes" in slide_data and slide_data["image_bytes"]
        
        # Set content area width based on whether there's an image
        if has_image:
            content_width = Inches(4.5)
            content_left = Inches(0.5)
        else:
            content_width = Inches(9)
            content_left = Inches(0.5)

        # Add Content (Bullet Points)
        content_box = slide.shapes.add_textbox(
            left=content_left,
            top=Inches(1.5),
            width=content_width,
            height=Inches(5.5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.margin_left = Inches(0.1)
        content_frame.margin_top = Inches(0.1)
        
        # Clear default paragraph
        content_frame.clear()
        
        # Add bullet points
        for i, point in enumerate(slide_data.get("content", [])):
            p = content_frame.add_paragraph()
            p.text = point
            p.font.size = Pt(16)
            p.level = 0
            p.space_before = Pt(6)
            p.space_after = Pt(6)
            p.line_spacing = 1.2

        # Add Image if available (right side, properly fitted)
        if has_image:
            img_stream = BytesIO(slide_data["image_bytes"])
            
            # Image positioning with proper margins
            img_left = Inches(5.3)
            img_top = Inches(1.5)  # Start below title
            img_max_width = Inches(4.2)
            img_max_height = Inches(5.7)  # Fit between title and bottom margin
            
            # Add picture and let it scale proportionally
            pic = slide.shapes.add_picture(
                img_stream, 
                img_left, 
                img_top, 
                width=img_max_width
            )
            
            # If height exceeds max, scale down to fit
            if pic.height > img_max_height:
                # Calculate aspect ratio
                aspect_ratio = pic.width / pic.height
                pic.height = img_max_height
                pic.width = int(img_max_height * aspect_ratio)
                
                # Center horizontally in the available space
                available_width = Inches(4.2)
                if pic.width < available_width:
                    pic.left = img_left + (available_width - pic.width) // 2

    output_path = os.path.join("backend", "generated", filename)
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    return output_path
